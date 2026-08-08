#!/usr/bin/env python3
"""
pptx_builder.py — Assignment OS PPTX Generator
Usage: python pptx_builder.py --draft <path> --output <path> [options]
Reads draft.json, generates presentation.pptx using python-pptx.
"""

import argparse
import json
import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ─── Color utilities ───────────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> tuple:
    h = hex_color.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def luminance(hex_color: str) -> float:
    r, g, b = hex_to_rgb(hex_color)
    def ch(c):
        c = c / 255
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def contrast_ratio(fg: str, bg: str) -> float:
    l1, l2 = luminance(fg), luminance(bg)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def choose_text_color(bg_hex: str) -> str:
    """Return #FFFFFF or #1B1B1B based on contrast against background."""
    white_ratio = contrast_ratio('#FFFFFF', bg_hex)
    dark_ratio = contrast_ratio('#1B1B1B', bg_hex)
    return '#FFFFFF' if white_ratio >= dark_ratio else '#1B1B1B'


# ─── Slide building helpers ─────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, hex_color: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    r, g, b = hex_to_rgb(hex_color)
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_text_box(slide, left_in, top_in, width_in, height_in, text,
                 font_pt, bold=False, color_hex='#FFFFFF',
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    r, g, b = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(r, g, b)
    return txBox


def add_bullet_box(slide, left_in, top_in, width_in, height_in, bullets,
                   font_pt, text_color='#FFFFFF', accent_hex='#E85D04'):
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tr, tg, tb = hex_to_rgb(text_color)
    ar, ag, ab = hex_to_rgb(accent_hex)
    for i, bullet in enumerate(bullets[:5]):  # max 5 bullets
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"→  {bullet}"
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(tr, tg, tb)
    return txBox


def get_heading(section: dict) -> str:
    return section.get('action_title') or section.get('heading', '')


def get_bullets(section: dict) -> list:
    if section.get('bullets'):
        return section['bullets']
    content = section.get('content', '')
    if isinstance(content, list):
        return content
    return content.split('\n') if content else []


def add_speaker_notes(slide, notes_text: str):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ─── Slide type builders ────────────────────────────────────────────

def build_title_cover(prs, section, brand_primary, brand_bg):
    slide = prs.slides.add_slide(prs.slide_master.slide_layouts[6])  # Blank
    set_slide_bg(slide, brand_bg)
    text_color = choose_text_color(brand_bg)
    # Accent bar top
    left = Inches(0); top = Inches(0); width = Inches(0.5); height = Inches(7.5)
    bar = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    bar.fill.solid()
    r, g, b = hex_to_rgb(brand_primary)
    bar.fill.fore_color.rgb = RGBColor(r, g, b)
    bar.line.fill.background()
    # Company label
    add_text_box(slide, 0.8, 0.4, 12, 0.5,
                 section.get('company', '').upper(), 11, bold=False,
                 color_hex=brand_primary, align=PP_ALIGN.LEFT)
    # Title
    add_text_box(slide, 0.8, 2.5, 11, 2,
                 get_heading(section), 36, bold=True,
                 color_hex=text_color, align=PP_ALIGN.LEFT)
    # Meta
    meta = section.get('meta', '')
    add_text_box(slide, 0.8, 6.8, 11, 0.5,
                 meta, 11, bold=False, color_hex=brand_primary, align=PP_ALIGN.LEFT)
    add_speaker_notes(slide, section.get('slide_notes', ''))
    return slide


def build_insight(prs, section, brand_primary, brand_bg):
    slide = prs.slides.add_slide(prs.slide_master.slide_layouts[6])
    set_slide_bg(slide, brand_bg)
    text_color = choose_text_color(brand_bg)
    # Heading
    add_text_box(slide, 0.5, 0.3, 12.3, 1.2,
                 get_heading(section), 24, bold=True, color_hex=text_color)
    # Bullets
    bullets = get_bullets(section)
    add_bullet_box(slide, 0.5, 1.8, 8, 4.5, bullets, 14,
                   text_color=text_color, accent_hex=brand_primary)
    # Stat callout
    stat = section.get('stat', '')
    stat_label = section.get('stat_label', '')
    if stat:
        callout = slide.shapes.add_shape(1, Inches(9.2), Inches(1.8), Inches(3.8), Inches(2.5))
        callout.fill.solid()
        r, g, b = hex_to_rgb(brand_primary)
        callout.fill.fore_color.rgb = RGBColor(r, g, b)
        callout.line.fill.background()
        add_text_box(slide, 9.4, 2.0, 3.4, 1.0,
                     stat, 32, bold=True, color_hex='#FFFFFF', align=PP_ALIGN.CENTER)
        add_text_box(slide, 9.4, 3.0, 3.4, 0.6,
                     stat_label, 10, color_hex='#FFFFFF', align=PP_ALIGN.CENTER)
    # Footnote
    footnote = section.get('footnote', '')
    if footnote:
        add_text_box(slide, 0.5, 7.0, 12.3, 0.35, footnote, 9, color_hex='#888888')
    add_speaker_notes(slide, section.get('slide_notes', ''))
    return slide


def build_recommendation(prs, section, brand_primary, brand_bg):
    slide = prs.slides.add_slide(prs.slide_master.slide_layouts[6])
    set_slide_bg(slide, brand_bg)
    text_color = choose_text_color(brand_bg)
    # Tag
    add_text_box(slide, 0.5, 0.2, 3, 0.35,
                 'RECOMMENDATION', 9, bold=True, color_hex=brand_primary)
    # Heading
    add_text_box(slide, 0.5, 0.5, 12.3, 1.5,
                 get_heading(section), 26, bold=True, color_hex=text_color)
    # Bullets
    bullets = get_bullets(section)
    add_bullet_box(slide, 0.5, 2.2, 12.3, 3.5, bullets, 14,
                   text_color=text_color, accent_hex=brand_primary)
    # Impact strip
    impact = section.get('impact', '')
    if impact:
        strip = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
        strip.fill.solid()
        r, g, b = hex_to_rgb(brand_primary)
        strip.fill.fore_color.rgb = RGBColor(r, g, b)
        strip.line.fill.background()
        add_text_box(slide, 0.3, 6.85, 12.7, 0.5,
                     f"Expected impact: {impact}", 12, bold=True, color_hex='#FFFFFF')
    add_speaker_notes(slide, section.get('slide_notes', ''))
    return slide


def build_generic(prs, section, brand_primary, brand_bg):
    """Fallback builder for slide types without dedicated handlers."""
    slide = prs.slides.add_slide(prs.slide_master.slide_layouts[6])
    set_slide_bg(slide, brand_bg)
    text_color = choose_text_color(brand_bg)
    # Heading
    add_text_box(slide, 0.5, 0.3, 12.3, 1.3,
                 get_heading(section), 24, bold=True, color_hex=text_color)
    # Content
    bullets = get_bullets(section)
    if bullets:
        add_bullet_box(slide, 0.5, 1.9, 12.3, 4.8, bullets, 14,
                       text_color=text_color, accent_hex=brand_primary)
    add_speaker_notes(slide, section.get('slide_notes', ''))
    return slide


BUILDERS = {
    'title_cover': build_title_cover,
    'executive_summary': build_insight,
    'insight': build_insight,
    'recommendation': build_recommendation,
    'comparison': build_generic,
    'metrics': build_generic,
    'timeline': build_generic,
    'process': build_generic,
    'section_divider': build_generic,
    'quote': build_generic,
    'two_column': build_generic,
    'appendix': build_generic,
}


# ─── Main ───────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='Build PPTX from draft.json')
    parser.add_argument('--draft', required=True, help='Path to draft.json')
    parser.add_argument('--output', required=True, help='Output .pptx path')
    parser.add_argument('--template', default=None, help='Base template .pptx (optional)')
    parser.add_argument('--brand-primary', default='#E85D04', help='Brand primary hex')
    parser.add_argument('--brand-bg', default='#0D1B2A', help='Brand background hex')
    parser.add_argument('--company', default='', help='Company name for the title slide')
    args = parser.parse_args()

    # Load draft
    with open(args.draft, 'r', encoding='utf-8') as f:
        draft = json.load(f)

    brand_primary = args.brand_primary
    brand_bg = args.brand_bg

    # Contrast check
    ratio = contrast_ratio('#FFFFFF', brand_bg)
    if ratio < 4.5:
        print(f"WARNING: White text on {brand_bg} has contrast {ratio:.1f}:1 (WCAG AA needs 4.5:1)")

    # Init presentation
    if args.template and Path(args.template).exists():
        prs = Presentation(args.template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    # Build slides
    all_sections = draft.get('sections', []) + draft.get('appendix_sections', [])
    slide_count = 0

    title_section = {
        'company': args.company,
        'action_title': draft.get('title', ''),
        'meta': draft.get('subtitle', ''),
    }
    build_title_cover(prs, title_section, brand_primary, brand_bg)
    slide_count += 1

    for section in all_sections:
        slide_type = section.get('type', 'insight')
        builder = BUILDERS.get(slide_type, build_generic)
        builder(prs, section, brand_primary, brand_bg)
        slide_count += 1

    # Save
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"PPTX Builder complete.")
    print(f"Slides: {slide_count}")
    print(f"Brand color: {brand_primary}")
    print(f"File: {output_path}")

    # Quick QA
    prs2 = Presentation(str(output_path))
    placeholders_found = 0
    for slide in prs2.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                for bad in ['[INSERT]', '[TODO]', 'PLACEHOLDER', 'lorem ipsum']:
                    if bad.lower() in text.lower():
                        placeholders_found += 1
                        print(f"  PLACEHOLDER WARNING: '{text[:60]}'")
    print(f"Placeholders found: {placeholders_found}")


def selftest():
    """python3 pptx_builder.py --selftest — guards the schema coupling to draft.json.

    This script was written against a schema case-builder never emitted (`action_title`,
    `bullets`), so it silently rendered blank headings and crashed on `content` being a
    list. Nothing caught it. These asserts fail the moment the reader and the real
    draft.json disagree again."""
    real = {"heading": "New investors stall at three points.",
            "content": ["Zero-candidate isn't one problem.", "It's three."]}
    assert get_heading(real) == "New investors stall at three points."
    assert get_bullets(real) == real["content"]          # list content must pass through
    # Legacy/alternate shape still works, so an older draft.json doesn't break the build.
    assert get_heading({"action_title": "X"}) == "X"
    assert get_bullets({"bullets": ["a"]}) == ["a"]
    assert get_bullets({"content": "one\ntwo"}) == ["one", "two"]   # str content splits
    # Absent fields degrade to empty, never raise — a missing section must not kill a deck.
    assert get_heading({}) == "" and get_bullets({}) == []
    print("pptx_builder selftest: ok")


if __name__ == '__main__':
    if '--selftest' in sys.argv:
        selftest()
    else:
        main()
